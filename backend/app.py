"""
Kairos Smart Student Planner - Python Backend
Flask API for PDF processing, Word/Excel/PowerPoint import/export.

Run: python app.py
Requires: pip install -r requirements.txt
Optional (PDF-to-image): Install Poppler and add to PATH
Optional (Word/PPT-to-PDF): Install LibreOffice OR Microsoft Word (Windows)
"""

import io
import os
import re
import json
import zipfile
import tempfile
import subprocess
from pathlib import Path
import random
import time
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

from flask import Flask, request, send_file, send_from_directory, jsonify
from flask_cors import CORS

# PDF
from pypdf import PdfWriter, PdfReader

# Watermark generation
from reportlab.pdfgen import canvas as rl_canvas
from reportlab.lib.colors import Color

# Office formats
from docx import Document as DocxDocument
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH

import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN

# Optional: pdf2image (requires poppler)
try:
    from pdf2image import convert_from_bytes
    PDF2IMAGE_AVAILABLE = True
except ImportError:
    PDF2IMAGE_AVAILABLE = False

from PIL import Image

app = Flask(__name__)
CORS(app, resources={r"/api/*": {"origins": "*"}})

# ── FRONTEND STATIC FILE SERVING ─────────────────────────────────────────────
# Serve the HTML/CSS/JS frontend from the repo root (one directory up from backend/)
FRONTEND_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))

def _no_cache(resp):
    """Prevent any proxy/CDN/browser from caching this response."""
    resp.headers['Cache-Control'] = 'no-store, no-cache, must-revalidate, max-age=0'
    resp.headers['Pragma'] = 'no-cache'
    resp.headers['Expires'] = '0'
    return resp

@app.route('/')
def serve_index():
    return _no_cache(send_from_directory(FRONTEND_DIR, 'landing.html'))

@app.route('/google402260b3a598f3b9.html')
@app.route('/<path:dummy>/google402260b3a598f3b9.html')
def serve_google_verification(dummy=None):
    return "google-site-verification: google402260b3a598f3b9.html", 200, {'Content-Type': 'text/plain'}

@app.route('/<path:filename>')
def serve_frontend(filename):
    filepath = os.path.join(FRONTEND_DIR, filename)
    if os.path.isfile(filepath):
        resp = send_from_directory(FRONTEND_DIR, filename)
        # Service worker: no-cache + correct MIME (browsers reject SW otherwise)
        if filename == 'sw.js':
            resp.headers['Content-Type'] = 'application/javascript'
            resp.headers['Service-Worker-Allowed'] = '/'
            return _no_cache(resp)
        # Manifest: correct MIME
        if filename == 'manifest.json':
            resp.headers['Content-Type'] = 'application/manifest+json'
            return _no_cache(resp)
        # HTML pages: always no-cache so browsers pick up new deployments immediately
        if filename.endswith('.html'):
            return _no_cache(resp)
        # JS/CSS: short cache (1 hour) — version params handle busting
        if filename.endswith(('.js', '.css')):
            resp.headers['Cache-Control'] = 'public, max-age=3600'
        return resp
    # SPA fallback
    return _no_cache(send_from_directory(FRONTEND_DIR, 'landing.html'))


# ─────────────────────────────────────────────
#  HELPERS
# ─────────────────────────────────────────────

def _hex_to_rgb(hex_color: str):
    """Convert #RRGGBB or #RGB to (r, g, b) tuple."""
    h = hex_color.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)


def _apply_cell_style(cell, style: dict):
    """Apply a style dict (from the Excel editor) to an openpyxl cell."""
    bold = bool(style.get("bold"))
    italic = bool(style.get("italic"))
    underline = "single" if style.get("underline") else None
    font_color = style.get("fontColor", "")
    fill_color = style.get("fillColor", "")
    align_class = style.get("alignClass", "")

    font_kwargs = {"bold": bold, "italic": italic}
    if underline:
        font_kwargs["underline"] = underline
    if font_color and font_color.startswith("#"):
        r, g, b = _hex_to_rgb(font_color)
        font_kwargs["color"] = f"{r:02X}{g:02X}{b:02X}"
    cell.font = Font(**font_kwargs)

    if fill_color and fill_color.startswith("#"):
        r, g, b = _hex_to_rgb(fill_color)
        cell.fill = PatternFill(fill_type="solid", fgColor=f"{r:02X}{g:02X}{b:02X}")

    h_align = "center" if "center" in align_class else ("right" if "right" in align_class else "left")
    cell.alignment = Alignment(horizontal=h_align)


def _cell_style_to_dict(cell) -> dict:
    """Convert openpyxl cell style to a dict matching the Excel editor's format."""
    s = {}
    if cell.font:
        s["bold"] = bool(cell.font.bold)
        s["italic"] = bool(cell.font.italic)
        s["underline"] = cell.font.underline == "single"
        if cell.font.color and cell.font.color.type == "rgb":
            rgb = cell.font.color.rgb  # 8-char ARGB string
            if rgb and len(rgb) == 8:
                s["fontColor"] = f"#{rgb[2:]}"
    if cell.fill and cell.fill.patternType == "solid":
        rgb = cell.fill.fgColor.rgb
        if rgb and len(rgb) == 8 and rgb.upper() != "FF000000":
            s["fillColor"] = f"#{rgb[2:]}"
    return s


def _html_to_docx(doc, html: str):
    """Parse simple HTML from the Word editor and populate a python-docx Document."""
    from html.parser import HTMLParser

    class _Builder(HTMLParser):
        def __init__(self):
            super().__init__()
            self.para = None
            self.bold = False
            self.italic = False
            self.underline = False
            self.color = None
            self.font_size = None

        def _ensure_para(self):
            if self.para is None:
                self.para = doc.add_paragraph()

        def handle_starttag(self, tag, attrs):
            attrs_dict = dict(attrs)
            style = attrs_dict.get("style", "")

            if tag in ("p", "div"):
                self.para = doc.add_paragraph()
                if "text-align:center" in style.replace(" ", "") or "text-align: center" in style:
                    self.para.alignment = WD_ALIGN_PARAGRAPH.CENTER
                elif "text-align:right" in style.replace(" ", "") or "text-align: right" in style:
                    self.para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                elif "text-align:justify" in style.replace(" ", ""):
                    self.para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
                else:
                    self.para.alignment = WD_ALIGN_PARAGRAPH.LEFT

            elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(tag[1])
                self.para = doc.add_heading("", level=min(level, 4))

            elif tag == "br":
                self._ensure_para()
                self.para.add_run("\n")

            elif tag in ("b", "strong"):
                self.bold = True
            elif tag in ("i", "em"):
                self.italic = True
            elif tag == "u":
                self.underline = True

            elif tag == "span":
                m = re.search(r"color:\s*([#\w()%,. ]+?)(?:;|$)", style)
                if m:
                    self.color = m.group(1).strip()
                m = re.search(r"font-size:\s*(\d+)", style)
                if m:
                    self.font_size = int(m.group(1))

        def handle_endtag(self, tag):
            if tag in ("p", "div", "h1", "h2", "h3", "h4", "h5", "h6"):
                self.para = None
            elif tag in ("b", "strong"):
                self.bold = False
            elif tag in ("i", "em"):
                self.italic = False
            elif tag == "u":
                self.underline = False
            elif tag == "span":
                self.color = None
                self.font_size = None

        def handle_data(self, data):
            if not data:
                return
            self._ensure_para()
            run = self.para.add_run(data)
            run.bold = self.bold
            run.italic = self.italic
            run.underline = self.underline
            if self.font_size:
                run.font.size = Pt(self.font_size)
            if self.color and self.color.startswith("#"):
                try:
                    r, g, b = _hex_to_rgb(self.color)
                    run.font.color.rgb = RGBColor(r, g, b)
                except Exception:
                    pass

    _Builder().feed(html)


def _docx_to_html(doc) -> str:
    """Convert a python-docx Document to HTML compatible with the Word editor."""
    parts = []
    align_map = {
        WD_ALIGN_PARAGRAPH.CENTER: "center",
        WD_ALIGN_PARAGRAPH.RIGHT: "right",
        WD_ALIGN_PARAGRAPH.JUSTIFY: "justify",
    }

    for para in doc.paragraphs:
        if not para.text and not para.runs:
            parts.append("<p><br></p>")
            continue

        align = align_map.get(para.alignment, "left")

        # Detect heading style
        tag = "p"
        if para.style and para.style.name.startswith("Heading"):
            try:
                lvl = int(para.style.name.split()[-1])
                tag = f"h{min(lvl, 4)}"
            except (ValueError, IndexError):
                pass

        inner = ""
        for run in para.runs:
            text = (run.text
                    .replace("&", "&amp;")
                    .replace("<", "&lt;")
                    .replace(">", "&gt;"))
            span_style = ""
            if run.font.size:
                span_style += f"font-size:{int(run.font.size.pt)}px;"
            if run.font.color and run.font.color.type:
                try:
                    rgb = run.font.color.rgb
                    span_style += f"color:#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x};"
                except Exception:
                    pass
            if span_style:
                text = f'<span style="{span_style}">{text}</span>'
            if run.bold:
                text = f"<b>{text}</b>"
            if run.italic:
                text = f"<i>{text}</i>"
            if run.underline:
                text = f"<u>{text}</u>"
            inner += text

        parts.append(f'<{tag} style="text-align:{align}">{inner}</{tag}>')

    return "\n".join(parts)


def _add_pptx_element(slide, el: dict, prs):
    """Place a slide element (from the editor's JSON) onto a python-pptx slide."""
    el_type = el.get("type", "body")
    text = el.get("text", "")

    # Map editor pixel coordinates (800×450) to slide EMU dimensions
    sw = prs.slide_width.inches
    sh = prs.slide_height.inches
    EDITOR_W, EDITOR_H = 800, 450

    x = PptInches(el.get("x", 0) / EDITOR_W * sw)
    y = PptInches(el.get("y", 0) / EDITOR_H * sh)
    w = PptInches(max(el.get("w", 200) / EDITOR_W * sw, 0.3))
    h = PptInches(max(el.get("h", 60) / EDITOR_H * sh, 0.2))

    # Shapes (rectangles, ovals, arrows, dividers)
    if el_type == "shape":
        SHAPE_IDS = {"rect": 1, "oval": 9, "arrow": 13}
        shape_id = SHAPE_IDS.get(el.get("shape", "rect"), 1)
        shape = slide.shapes.add_shape(shape_id, x, y, w, h)
        fill_color = el.get("fill", "#c44529")
        if fill_color and fill_color != "transparent":
            try:
                r, g, b = _hex_to_rgb(fill_color)
                shape.fill.solid()
                shape.fill.fore_color.rgb = PptRGBColor(r, g, b)
            except Exception:
                pass
        shape.line.fill.background()
        if text:
            shape.text_frame.text = text
        return

    # Text / title / subtitle / body
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    font_size = el.get("fontSize", 14)
    bold = str(el.get("fontWeight", "400")) in ("700", "bold")
    color = el.get("color", "#333333")
    text_align = el.get("textAlign", "left")
    ALIGN_MAP = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}

    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ALIGN_MAP.get(text_align, PP_ALIGN.LEFT)
        run = p.add_run()
        run.text = line
        run.font.size = PptPt(font_size)
        run.font.bold = bold
        if color and color != "transparent":
            try:
                r, g, b = _hex_to_rgb(color)
                run.font.color.rgb = PptRGBColor(r, g, b)
            except Exception:
                pass


def _convert_with_libreoffice(input_path: str, out_dir: str) -> str | None:
    """Try to convert a file to PDF using LibreOffice. Returns output path or None."""
    for cmd in ["soffice", "libreoffice"]:
        try:
            result = subprocess.run(
                [cmd, "--headless", "--convert-to", "pdf", "--outdir", out_dir, input_path],
                capture_output=True, timeout=60
            )
            pdf_name = Path(input_path).stem + ".pdf"
            pdf_path = os.path.join(out_dir, pdf_name)
            if os.path.exists(pdf_path):
                return pdf_path
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None


# ─────────────────────────────────────────────
#  PDF ROUTES
# ─────────────────────────────────────────────

@app.route("/api/pdf/merge", methods=["POST"])
def pdf_merge():
    files = request.files.getlist("files")
    if len(files) < 2:
        return jsonify({"error": "Upload at least 2 PDF files"}), 400

    writer = PdfWriter()
    for f in files:
        reader = PdfReader(f.stream)
        for page in reader.pages:
            writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, mimetype="application/pdf",
                     as_attachment=True, download_name="merged.pdf")


@app.route("/api/pdf/split", methods=["POST"])
def pdf_split():
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    reader = PdfReader(f.stream)
    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            page_buf = io.BytesIO()
            writer.write(page_buf)
            zf.writestr(f"page_{i + 1}.pdf", page_buf.getvalue())

    zip_buf.seek(0)
    return send_file(zip_buf, mimetype="application/zip",
                     as_attachment=True, download_name="split_pages.zip")


@app.route("/api/pdf/compress", methods=["POST"])
def pdf_compress():
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    reader = PdfReader(f.stream)
    writer = PdfWriter()
    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)
    writer.compress_identical_objects(remove_identicals=True, remove_orphans=True)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, mimetype="application/pdf",
                     as_attachment=True, download_name="compressed.pdf")


@app.route("/api/pdf/watermark", methods=["POST"])
def pdf_watermark():
    f = request.files.get("file")
    text = request.form.get("text", "WATERMARK")
    opacity = float(request.form.get("opacity", "0.3"))
    if not f:
        return jsonify({"error": "No file provided"}), 400

    reader = PdfReader(f.stream)
    writer = PdfWriter()

    for page in reader.pages:
        pw = float(page.mediabox.width)
        ph = float(page.mediabox.height)

        # Build a per-page watermark at the exact page size
        wm_buf = io.BytesIO()
        c = rl_canvas.Canvas(wm_buf, pagesize=(pw, ph))
        c.saveState()
        c.setFont("Helvetica-Bold", max(int(min(pw, ph) / 8), 20))
        c.setFillColor(Color(0.5, 0.5, 0.5, alpha=opacity))
        c.translate(pw / 2, ph / 2)
        c.rotate(45)
        c.drawCentredString(0, 0, text)
        c.restoreState()
        c.save()
        wm_buf.seek(0)

        wm_page = PdfReader(wm_buf).pages[0]
        page.merge_page(wm_page)
        writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, mimetype="application/pdf",
                     as_attachment=True, download_name="watermarked.pdf")


@app.route("/api/pdf/to-image", methods=["POST"])
def pdf_to_image():
    if not PDF2IMAGE_AVAILABLE:
        return jsonify({
            "error": "pdf2image not available. Install poppler and run: pip install pdf2image"
        }), 500

    f = request.files.get("file")
    fmt = request.form.get("format", "jpg").lower()
    if not f:
        return jsonify({"error": "No file provided"}), 400

    pdf_bytes = f.read()
    images = convert_from_bytes(pdf_bytes, dpi=150)
    pil_fmt = "JPEG" if fmt == "jpg" else "PNG"
    mime = "image/jpeg" if fmt == "jpg" else "image/png"

    if len(images) == 1:
        buf = io.BytesIO()
        img = images[0].convert("RGB") if fmt == "jpg" else images[0]
        img.save(buf, format=pil_fmt)
        buf.seek(0)
        return send_file(buf, mimetype=mime, as_attachment=True,
                         download_name=f"page_1.{fmt}")

    zip_buf = io.BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for i, img in enumerate(images):
            buf = io.BytesIO()
            img = img.convert("RGB") if fmt == "jpg" else img
            img.save(buf, format=pil_fmt)
            zf.writestr(f"page_{i + 1}.{fmt}", buf.getvalue())
    zip_buf.seek(0)
    return send_file(zip_buf, mimetype="application/zip",
                     as_attachment=True, download_name=f"pages.zip")


@app.route("/api/pdf/rotate", methods=["POST"])
def pdf_rotate():
    f = request.files.get("file")
    angle = int(request.form.get("angle", "90"))
    if not f:
        return jsonify({"error": "No file provided"}), 400
    if angle not in (90, 180, 270):
        return jsonify({"error": "Angle must be 90, 180, or 270"}), 400

    reader = PdfReader(f.stream)
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(angle)
        writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, mimetype="application/pdf",
                     as_attachment=True, download_name="rotated.pdf")


@app.route("/api/pdf/unlock", methods=["POST"])
def pdf_unlock():
    f = request.files.get("file")
    password = request.form.get("password", "")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    reader = PdfReader(f.stream)
    if reader.is_encrypted:
        if not reader.decrypt(password):
            return jsonify({"error": "Incorrect password"}), 400

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, mimetype="application/pdf",
                     as_attachment=True, download_name="unlocked.pdf")


@app.route("/api/pdf/extract-pages", methods=["POST"])
def pdf_extract_pages():
    f = request.files.get("file")
    page_range = request.form.get("pages", "1")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    reader = PdfReader(f.stream)
    total = len(reader.pages)

    indices = set()
    for part in page_range.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            for p in range(int(a), int(b) + 1):
                if 1 <= p <= total:
                    indices.add(p - 1)
        elif part.isdigit():
            p = int(part)
            if 1 <= p <= total:
                indices.add(p - 1)

    if not indices:
        return jsonify({"error": "No valid pages in range"}), 400

    writer = PdfWriter()
    for i in sorted(indices):
        writer.add_page(reader.pages[i])

    out = io.BytesIO()
    writer.write(out)
    out.seek(0)
    return send_file(out, mimetype="application/pdf",
                     as_attachment=True, download_name="extracted.pdf")


@app.route("/api/pdf/word-to-pdf", methods=["POST"])
def word_to_pdf():
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, f.filename or "document.docx")
        f.save(input_path)

        # Try LibreOffice (cross-platform)
        pdf_path = _convert_with_libreoffice(input_path, tmpdir)
        if pdf_path:
            data = open(pdf_path, "rb").read()
            return send_file(io.BytesIO(data), mimetype="application/pdf",
                             as_attachment=True,
                             download_name=Path(f.filename).stem + ".pdf")

        # Fallback: docx2pdf (Windows with MS Word installed)
        try:
            from docx2pdf import convert
            out_path = os.path.join(tmpdir, Path(input_path).stem + ".pdf")
            convert(input_path, out_path)
            data = open(out_path, "rb").read()
            return send_file(io.BytesIO(data), mimetype="application/pdf",
                             as_attachment=True,
                             download_name=Path(f.filename).stem + ".pdf")
        except Exception as e:
            return jsonify({
                "error": (
                    "Conversion failed. Install LibreOffice (recommended) "
                    f"or Microsoft Word to enable this feature. Detail: {e}"
                )
            }), 500


@app.route("/api/pdf/ppt-to-pdf", methods=["POST"])
def ppt_to_pdf():
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, f.filename or "presentation.pptx")
        f.save(input_path)

        pdf_path = _convert_with_libreoffice(input_path, tmpdir)
        if pdf_path:
            data = open(pdf_path, "rb").read()
            return send_file(io.BytesIO(data), mimetype="application/pdf",
                             as_attachment=True,
                             download_name=Path(f.filename).stem + ".pdf")

        try:
            from docx2pdf import convert
            out_path = os.path.join(tmpdir, Path(input_path).stem + ".pdf")
            convert(input_path, out_path)
            data = open(out_path, "rb").read()
            return send_file(io.BytesIO(data), mimetype="application/pdf",
                             as_attachment=True,
                             download_name=Path(f.filename).stem + ".pdf")
        except Exception as e:
            return jsonify({
                "error": (
                    "Conversion failed. Install LibreOffice to enable this feature. "
                    f"Detail: {e}"
                )
            }), 500


# ─────────────────────────────────────────────
#  WORD ROUTES
# ─────────────────────────────────────────────

@app.route("/api/word/export", methods=["POST"])
def word_export():
    """Export the Word editor's HTML content as a real .docx file."""
    data = request.get_json(force=True)
    html_content = data.get("html", "")
    doc_name = data.get("name", "Document")

    doc = DocxDocument()
    # Remove default empty paragraph added by python-docx
    for p in list(doc.paragraphs):
        p._element.getparent().remove(p._element)

    _html_to_docx(doc, html_content)

    out = io.BytesIO()
    doc.save(out)
    out.seek(0)
    fname = doc_name if doc_name.endswith(".docx") else doc_name + ".docx"
    return send_file(
        out,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        as_attachment=True,
        download_name=fname,
    )


@app.route("/api/word/import", methods=["POST"])
def word_import():
    """Import a .docx file and return HTML for the Word editor."""
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    doc = DocxDocument(f.stream)
    html = _docx_to_html(doc)
    return jsonify({"html": html, "name": Path(f.filename).stem})


# ─────────────────────────────────────────────
#  EXCEL ROUTES
# ─────────────────────────────────────────────

@app.route("/api/excel/export", methods=["POST"])
def excel_export():
    """Export the Excel editor's sheet data as a real .xlsx file."""
    data = request.get_json(force=True)
    sheets_data = data.get("sheets", [])
    workbook_name = data.get("name", "Workbook")

    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    for sheet in sheets_data:
        ws = wb.create_sheet(title=sheet.get("name", "Sheet"))
        grid = sheet.get("data", [])
        styles = sheet.get("styles", [])
        col_widths = sheet.get("colWidths", [])

        for r_i, row in enumerate(grid):
            for c_i, val in enumerate(row):
                if val not in (None, ""):
                    cell = ws.cell(row=r_i + 1, column=c_i + 1, value=val)
                    if r_i < len(styles) and c_i < len(styles[r_i]):
                        s = styles[r_i][c_i]
                        if s:
                            _apply_cell_style(cell, s)

        for c_i, px_width in enumerate(col_widths):
            col_letter = openpyxl.utils.get_column_letter(c_i + 1)
            ws.column_dimensions[col_letter].width = max(px_width / 7.5, 8)

    out = io.BytesIO()
    wb.save(out)
    out.seek(0)
    fname = workbook_name if workbook_name.endswith(".xlsx") else workbook_name + ".xlsx"
    return send_file(
        out,
        mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        as_attachment=True,
        download_name=fname,
    )


@app.route("/api/excel/import", methods=["POST"])
def excel_import():
    """Import a .xlsx/.csv file and return sheet data for the Excel editor."""
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    ROWS, COLS = 50, 26
    fname = f.filename or ""

    # CSV import
    if fname.lower().endswith(".csv"):
        import csv
        content = f.read().decode("utf-8-sig", errors="replace")
        reader = csv.reader(content.splitlines())
        grid = [[""] * COLS for _ in range(ROWS)]
        for r_i, row in enumerate(reader):
            if r_i >= ROWS:
                break
            for c_i, val in enumerate(row):
                if c_i < COLS:
                    grid[r_i][c_i] = val
        return jsonify({
            "sheets": [{
                "name": Path(fname).stem,
                "data": grid,
                "styles": [[{} for _ in range(COLS)] for _ in range(ROWS)],
                "colWidths": [80] * COLS,
                "rowHeights": [22] * ROWS,
            }],
            "name": Path(fname).stem,
        })

    # XLSX import
    wb = openpyxl.load_workbook(f.stream, data_only=True)
    sheets_out = []
    for ws in wb.worksheets:
        grid = [[""] * COLS for _ in range(ROWS)]
        styles = [[{} for _ in range(COLS)] for _ in range(ROWS)]
        for row in ws.iter_rows():
            for cell in row:
                r, c = cell.row - 1, cell.column - 1
                if r < ROWS and c < COLS:
                    v = cell.value
                    grid[r][c] = str(v) if v is not None else ""
                    styles[r][c] = _cell_style_to_dict(cell)
        col_widths = [80] * COLS
        for i in range(COLS):
            letter = openpyxl.utils.get_column_letter(i + 1)
            if letter in ws.column_dimensions:
                w = ws.column_dimensions[letter].width or 8
                col_widths[i] = int(w * 7.5)
        sheets_out.append({
            "name": ws.title,
            "data": grid,
            "styles": styles,
            "colWidths": col_widths,
            "rowHeights": [22] * ROWS,
        })

    return jsonify({"sheets": sheets_out, "name": Path(fname).stem})


# ─────────────────────────────────────────────
#  POWERPOINT ROUTES
# ─────────────────────────────────────────────

@app.route("/api/ppt/export", methods=["POST"])
def ppt_export():
    """Export the PowerPoint editor's slide data as a real .pptx file."""
    data = request.get_json(force=True)
    slides_data = data.get("slides", [])
    pres_name = data.get("name", "Presentation")

    prs = Presentation()
    # Widescreen 16:9
    prs.slide_width = Emu(9144000)   # 10 inches
    prs.slide_height = Emu(5143500)  # 5.625 inches

    blank_layout = prs.slide_layouts[6]

    for sd in slides_data:
        slide = prs.slides.add_slide(blank_layout)

        # Background color
        bg_hex = sd.get("bg", "#ffffff")
        try:
            r, g, b = _hex_to_rgb(bg_hex)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = PptRGBColor(r, g, b)
        except Exception:
            pass

        for el in sd.get("elements", []):
            try:
                _add_pptx_element(slide, el, prs)
            except Exception:
                pass

        notes_text = sd.get("notes", "")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    out = io.BytesIO()
    prs.save(out)
    out.seek(0)
    fname = pres_name if pres_name.endswith(".pptx") else pres_name + ".pptx"
    return send_file(
        out,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=fname,
    )


@app.route("/api/ppt/import", methods=["POST"])
def ppt_import():
    """Import a .pptx file and return slide data for the PowerPoint editor."""
    f = request.files.get("file")
    if not f:
        return jsonify({"error": "No file provided"}), 400

    prs = Presentation(f.stream)
    EDITOR_W, EDITOR_H = 800, 450
    sw = prs.slide_width.pt or 720
    sh = prs.slide_height.pt or 540

    slides_out = []
    for slide in prs.slides:
        elements = []
        bg_color = "#ffffff"

        try:
            bg_fill = slide.background.fill
            if hasattr(bg_fill, "fore_color") and bg_fill.fore_color.type:
                rgb = bg_fill.fore_color.rgb
                bg_color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
        except Exception:
            pass

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            x = int(shape.left.pt / sw * EDITOR_W) if sw else 0
            y = int(shape.top.pt / sh * EDITOR_H) if sh else 0
            w = int(shape.width.pt / sw * EDITOR_W) if sw else 200
            h = int(shape.height.pt / sh * EDITOR_H) if sh else 80

            text = shape.text_frame.text
            fs, fw, color, align = 14, "400", "#333", "left"

            for para in shape.text_frame.paragraphs:
                if para.runs:
                    run = para.runs[0]
                    if run.font.size:
                        fs = int(run.font.size.pt)
                    if run.font.bold:
                        fw = "700"
                    try:
                        if run.font.color and run.font.color.type:
                            rgb = run.font.color.rgb
                            color = f"#{rgb.red:02x}{rgb.green:02x}{rgb.blue:02x}"
                    except Exception:
                        pass
                if para.alignment:
                    align = {1: "left", 2: "center", 3: "right"}.get(para.alignment, "left")
                break

            el_type = "body"
            name = shape.name.lower()
            if "title" in name:
                el_type = "title"
                fs = max(fs, 24)
                fw = "700"
            elif "subtitle" in name:
                el_type = "subtitle"

            elements.append({
                "id": f"_imp{shape.shape_id}",
                "type": el_type,
                "x": x, "y": y, "w": w, "h": h,
                "text": text,
                "fontSize": fs,
                "fontWeight": fw,
                "textAlign": align,
                "color": color,
            })

        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text
        except Exception:
            pass

        slides_out.append({
            "id": f"_s{len(slides_out)}",
            "elements": elements,
            "notes": notes_text,
            "bg": bg_color,
            "theme": "Default",
            "transition": "none",
            "layout": "title",
        })

    return jsonify({"slides": slides_out, "name": Path(f.filename).stem})


# ─────────────────────────────────────────────
#  AUTH ROUTES
# ─────────────────────────────────────────────

otp_store = {}

SMTP_SERVER = os.environ.get("SMTP_SERVER", "smtp.gmail.com")
SMTP_PORT = int(os.environ.get("SMTP_PORT", 587))
SMTP_EMAIL = os.environ.get("SMTP_EMAIL", "")
SMTP_PASSWORD = os.environ.get("SMTP_PASSWORD", "")

@app.route("/api/auth/send-otp", methods=["POST"])
def send_otp():
    data = request.get_json(force=True)
    email = data.get("email")
    if not email:
        return jsonify({"error": "Email is required"}), 400
    
    if not SMTP_EMAIL or not SMTP_PASSWORD:
        print(f"[AUTH] SMTP not configured. Generating dummy OTP for {email}")
        otp = "123456"
    else:
        otp = str(random.randint(100000, 999999))
    
    otp_store[email] = {
        "otp": otp,
        "expires_at": time.time() + 600 # 10 minutes expiry
    }
    
    if SMTP_EMAIL and SMTP_PASSWORD:
        try:
            msg = MIMEMultipart()
            msg['From'] = SMTP_EMAIL
            msg['To'] = email
            msg['Subject'] = "Your Kairos Verification Code"
            body = f"Hello,\n\nYour verification code for Kairos is: {otp}\n\nThis code will expire in 10 minutes.\n\nBest,\nKairos Team"
            msg.attach(MIMEText(body, 'plain'))
            
            server = smtplib.SMTP(SMTP_SERVER, SMTP_PORT)
            server.starttls()
            server.login(SMTP_EMAIL, SMTP_PASSWORD)
            server.send_message(msg)
            server.quit()
        except Exception as e:
            return jsonify({"error": f"Failed to send email: {str(e)}"}), 500
            
    return jsonify({"success": True, "message": "OTP sent successfully"})

@app.route("/api/auth/verify-otp", methods=["POST"])
def verify_otp():
    data = request.get_json(force=True)
    email = data.get("email")
    otp = data.get("otp")
    
    if not email or not otp:
        return jsonify({"error": "Email and OTP are required"}), 400
        
    record = otp_store.get(email)
    if not record:
        return jsonify({"error": "No OTP requested for this email"}), 400
        
    if time.time() > record["expires_at"]:
        del otp_store[email]
        return jsonify({"error": "OTP has expired"}), 400
        
    if record["otp"] != otp:
        return jsonify({"error": "Invalid OTP"}), 400
        
    # Verification successful
    del otp_store[email]
    return jsonify({"success": True, "message": "Email verified successfully"})

# ─────────────────────────────────────────────
#  HEALTH CHECK
# ─────────────────────────────────────────────

@app.route("/api/health", methods=["GET"])
def health():
    return jsonify({
        "status": "ok",
        "pdf2image": PDF2IMAGE_AVAILABLE,
    })


if __name__ == "__main__":
    print("=" * 55)
    print("  Kairos Backend Server")
    print("  http://localhost:5000")
    print("=" * 55)
    print(f"  pdf2image (PDF→image): {'YES' if PDF2IMAGE_AVAILABLE else 'NO (install poppler)'}")
    print("=" * 55)
    app.run(debug=False, host="0.0.0.0", port=5000)
